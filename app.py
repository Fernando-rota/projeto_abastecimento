import io
import tempfile
import pandas as pd
import streamlit as st
import plotly.express as px
from pptx import Presentation
from pptx.util import Inches

st.set_page_config(page_title="BI Abastecimento Frota + Export PPTX", layout="wide")
st.title("📊 BI Abastecimento - Interno e Externo + Exportação PPTX")

@st.cache_data
def load_data(file_path):
    interno = pd.read_excel(file_path, sheet_name='interno')
    externo = pd.read_excel(file_path, sheet_name='externo')

    for df in [interno, externo]:
        df.rename(columns=lambda x: x.strip().lower().replace(' ', '_'), inplace=True)

    interno['data'] = pd.to_datetime(interno['data'], errors='coerce')
    externo['data'] = pd.to_datetime(externo['data'], errors='coerce')

    interno.dropna(subset=['data'], inplace=True)
    externo.dropna(subset=['data'], inplace=True)

    return interno, externo

def preprocess_abastecimentos(df, litros_col, km_col, combust_col):
    df = df.rename(columns={
        litros_col: 'litros',
        km_col: 'km',
        combust_col: 'combustivel'
    })
    df['km'] = df['km'].astype(str).str.replace(',', '.', regex=False)
    df['litros'] = df['litros'].astype(str).str.replace(',', '.', regex=False)
    df['km'] = pd.to_numeric(df['km'], errors='coerce')
    df['litros'] = pd.to_numeric(df['litros'], errors='coerce')
    df = df.dropna(subset=['km', 'litros', 'placa'])
    return df[['data', 'placa', 'combustivel', 'litros', 'km']]

def fig_to_image(fig):
    tmp = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
    fig.write_image(tmp.name)
    return tmp.name

def criar_ppt(resumo_ab_df, fig_ab):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    # Slide 1: título
    slide1 = prs.slides.add_slide(blank_slide_layout)
    txBox = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    tf = txBox.text_frame
    tf.text = "Dashboard Abastecimento Frota - Resumo"

    # Slide 2: gráfico litros consumidos (aba interno + externo)
    slide2 = prs.slides.add_slide(blank_slide_layout)
    slide2.shapes.add_picture(fig_to_image(fig_ab), Inches(0.5), Inches(0.5), Inches(9), Inches(5))

    # Slide 3: tabela resumo aba interno + externo
    slide3 = prs.slides.add_slide(blank_slide_layout)
    text2 = "Indicadores Abastecimento Interno + Externo\n\n"
    for _, row in resumo_ab_df.iterrows():
        text2 += f"Placa: {row['placa']} | Total Litros: {row['total_litros']:.2f} | Média Km: {row['media_km']:.0f} | Registros: {row['registros']}\n"
    txBox3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6))
    tf3 = txBox3.text_frame
    tf3.text = text2

    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

uploaded_file = st.file_uploader("📁 Carregue sua planilha Excel com abas: interno e externo", type=['xlsx'])
if uploaded_file:
    interno, externo = load_data(uploaded_file)

    st.sidebar.header("Filtros Globais")

    placas_unicas = sorted(set(interno['placa'].dropna().unique()) | set(externo['placa'].dropna().unique()))
    placas_selected = st.sidebar.multiselect("Placas", placas_unicas, default=placas_unicas)

    combust_interno = interno['tipo'].dropna().unique() if 'tipo' in interno.columns else []
    combust_externo = externo['tipo_combustivel'].dropna().unique() if 'tipo_combustivel' in externo.columns else []
    combust_unificados = sorted(set(combust_interno) | set(combust_externo))
    combust_selected = st.sidebar.multiselect("Tipo Combustível", combust_unificados, default=combust_unificados)

    data_min = min(interno['data'].min(), externo['data'].min())
    data_max = max(interno['data'].max(), externo['data'].max())
    data_range = st.sidebar.date_input("Período", [data_min, data_max])

    data_start, data_end = pd.to_datetime(data_range[0]), pd.to_datetime(data_range[1])

    # Filtrar dados com base nos filtros
    interno_filt = interno[
        (interno['placa'].isin(placas_selected)) &
        (interno['data'] >= data_start) & (interno['data'] <= data_end)
    ]
    externo_filt = externo[
        (externo['placa'].isin(placas_selected)) &
        (externo['data'] >= data_start) & (externo['data'] <= data_end)
    ]

    if 'tipo' in interno_filt.columns:
        interno_filt = interno_filt[interno_filt['tipo'].isin(combust_selected)]
    if 'tipo_combustivel' in externo_filt.columns:
        externo_filt = externo_filt[externo_filt['tipo_combustivel'].isin(combust_selected)]

    # Processar para padrão unificado
    interno_proc = preprocess_abastecimentos(interno_filt, 'quantidade_de_litros', 'km_atual', 'tipo')
    externo_proc = preprocess_abastecimentos(externo_filt, 'quantidade_de_litros', 'km_atual', 'tipo_combustivel')

    abastecimentos = pd.concat([interno_proc, externo_proc], ignore_index=True)

    resumo_ab = abastecimentos.groupby('placa').agg(
        total_litros=('litros', 'sum'),
        media_km=('km', 'mean'),
        registros=('data', 'count')
    ).reset_index()

    resumo_ab = resumo_ab.sort_values('total_litros', ascending=False)

    # Mostrar indicadores
    st.header("⛽ Indicadores Abastecimento Interno + Externo")
    st.dataframe(resumo_ab.style.format({
        'total_litros': '{:,.2f}',
        'media_km': '{:,.0f}',
        'registros': '{:,.0f}'
    }))

    # Gráfico total litros por placa
    fig = px.bar(resumo_ab, x='placa', y='total_litros',
                 labels={'total_litros': 'Total Litros', 'placa': 'Placa'},
                 title='Total de Litros Consumidos por Veículo')
    st.plotly_chart(fig, use_container_width=True)

    # Botão exportar PPTX
    pptx_file = criar_ppt(resumo_ab, fig)
    st.download_button(
        label="📥 Exportar Apresentação PowerPoint",
        data=pptx_file,
        file_name="dashboard_abastecimento.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

else:
    st.info("Faça upload da planilha Excel para começar.")
